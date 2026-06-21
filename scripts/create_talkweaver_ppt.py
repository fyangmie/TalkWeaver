#!/usr/bin/env python3
"""Generate a concise TalkWeaver project overview deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "TalkWeaver_Project_Presentation.pptx"

INK = RGBColor(24, 32, 37)
MUTED = RGBColor(88, 101, 110)
TEAL = RGBColor(16, 121, 111)
AMBER = RGBColor(183, 121, 31)
RED = RGBColor(178, 59, 59)
BLUE = RGBColor(47, 111, 143)
PAPER = RGBColor(247, 250, 249)
LINE = RGBColor(216, 224, 228)
WHITE = RGBColor(255, 255, 255)


def set_run(run, *, size=18, bold=False, color=INK):
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, *, size=18, bold=False, color=INK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_bullets(slide, x, y, w, h, items, *, size=17, color=INK, gap=0.1):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(gap * 24)
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.62, 0.32, 9.9, 0.45, "TalkWeaver", size=13, bold=True, color=TEAL)
    add_textbox(slide, 0.62, 0.76, 10.8, 0.6, title, size=28, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, 0.64, 1.35, 11.0, 0.42, subtitle, size=14, color=MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(1.82), Inches(12.1), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_card(slide, x, y, w, h, title, body, *, accent=TEAL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    add_textbox(slide, x + 0.22, y + 0.14, w - 0.35, 0.32, title, size=15, bold=True, color=INK)
    add_textbox(slide, x + 0.22, y + 0.54, w - 0.35, h - 0.65, body, size=12.5, color=MUTED)


def add_metric(slide, x, y, w, label, value, note="", *, accent=TEAL):
    add_card(slide, x, y, w, 1.05, label, "", accent=accent)
    add_textbox(slide, x + 0.22, y + 0.34, w - 0.4, 0.38, value, size=19, bold=True, color=accent)
    if note:
        add_textbox(slide, x + 0.22, y + 0.72, w - 0.4, 0.24, note, size=10.5, color=MUTED)


def add_image(slide, path, x, y, w, h=None):
    source = ROOT / path
    if source.exists():
        if h is None:
            slide.shapes.add_picture(str(source), Inches(x), Inches(y), width=Inches(w))
        else:
            slide.shapes.add_picture(str(source), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        add_card(slide, x, y, w, h or 2.0, "缺少图表", str(path), accent=RED)


def add_footer(slide, num):
    add_textbox(slide, 11.95, 7.05, 0.7, 0.25, f"{num}/10", size=10.5, color=MUTED, align=PP_ALIGN.RIGHT)


def make_deck() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    add_textbox(slide, 0.72, 0.78, 11.4, 0.42, "AI Meeting Detective", size=16, bold=True, color=TEAL)
    add_textbox(slide, 0.72, 1.32, 11.5, 0.9, "TalkWeaver", size=52, bold=True, color=INK)
    add_textbox(slide, 0.78, 2.22, 11.4, 0.58, "把混乱多人语音变成可审计的对话证据地图", size=24, bold=True, color=INK)
    add_textbox(
        slide,
        0.8,
        3.12,
        10.6,
        0.8,
        "Evidence-grounded conversation maps for overlap, interruptions, misheard terms, and correction audits.",
        size=18,
        color=MUTED,
    )
    add_card(slide, 0.82, 4.4, 3.45, 1.0, "网页 Demo", "https://fyangmie-talkweaver.hf.space", accent=TEAL)
    add_card(slide, 4.55, 4.4, 3.45, 1.0, "主展示样本", "Earnings-22 多说话人术语修正", accent=BLUE)
    add_card(slide, 8.28, 4.4, 3.45, 1.0, "项目定位", "不是会议纪要系统，而是证据审计框架", accent=AMBER)
    add_footer(slide, 1)

    # 2. Motivation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "研究动机：会议转写不是只看 WER", "多人、重叠、低能量语音、专有词和 LLM 后处理会共同放大错误")
    add_bullets(
        slide,
        0.75,
        2.18,
        5.6,
        3.4,
        [
            "普通 ASR 输出通常是一段“看起来流畅”的文本，但用户不知道哪里不可靠。",
            "多人会议里，关键问题不是只识别文字，还包括谁说的、何时说的、是否重叠、哪些词被误听。",
            "LLM 后处理如果没有证据约束，可能把不确定内容改得更自然，却更不真实。",
            "TalkWeaver 的核心假设：混乱会议转写应被当作 evidence-grounded auditing problem。",
        ],
        size=18,
    )
    add_metric(slide, 7.0, 2.1, 2.7, "FLEURS zh-CN", "CER 0.113", "read speech, base", accent=TEAL)
    add_metric(slide, 10.0, 2.1, 2.7, "AISHELL-4", "CER 0.537", "meeting speech, base", accent=RED)
    add_metric(slide, 7.0, 3.55, 2.7, "AMI", "WER 0.398", "formal meeting, base", accent=AMBER)
    add_metric(slide, 10.0, 3.55, 2.7, "结论", "读句子 ≠ 会议", "需要证据地图", accent=BLUE)
    add_footer(slide, 2)

    # 3. Research questions and contributions
    slide = prs.slides.add_slide(blank)
    add_title(slide, "研究问题与贡献", "从“自动改文本”转向“带证据地组织、修正、审计文本”")
    add_card(slide, 0.75, 2.1, 3.0, 1.35, "RQ1", "说话人/时间结构能否让 transcript 更可审计？", accent=TEAL)
    add_card(slide, 3.95, 2.1, 3.0, 1.35, "RQ2", "overlap-aware uncertainty 能否降低不安全修正？", accent=AMBER)
    add_card(slide, 7.15, 2.1, 3.0, 1.35, "RQ3", "RAG glossary 能否恢复专业词和公司名？", accent=BLUE)
    add_card(slide, 10.35, 2.1, 2.25, 1.35, "RQ4", "本地 ASR 有何速度/质量权衡？", accent=RED)
    add_bullets(
        slide,
        0.92,
        4.1,
        11.4,
        1.8,
        [
            "贡献 1：Temporal-anchor ConversationMap，统一 ASR、speaker、time、overlap、retrieved terms、audit。",
            "贡献 2：Conservative RAG-based term recovery，只在证据足够时修正，弱证据进入 review。",
            "贡献 3：Safety analysis，区分真实公开数据、controlled stress-test、mock/demo 和负结果。",
        ],
        size=18,
    )
    add_footer(slide, 3)

    # 4. Architecture
    slide = prs.slides.add_slide(blank)
    add_title(slide, "系统架构：从音频到 EvidenceMap", "每一步都保留输入、输出和证据来源")
    add_image(slide, "assets/architecture.png", 0.8, 2.05, 5.45)
    add_bullets(
        slide,
        6.55,
        2.05,
        5.8,
        4.05,
        [
            "Audio preprocessing：统一采样率、声道和音量。",
            "ASR：faster-whisper / whisper.cpp / mock fallback。",
            "Diarization：pyannote 或 reference speaker-time evidence。",
            "Overlap & event extraction：标出重叠、打断候选、低置信度区域。",
            "RAG term rescue：检索专业词、公司名、财经术语。",
            "Correction audit：保留 raw、corrected、证据词和 accept/reject/review 决策。",
        ],
        size=16.5,
    )
    add_footer(slide, 4)

    # 5. EvidenceMap demo
    slide = prs.slides.add_slide(blank)
    add_title(slide, "网页 MVP：让用户看到“哪句话被改了”", "Demo 已部署到 Hugging Face Docker Space")
    add_card(
        slide,
        0.75,
        2.05,
        5.9,
        1.28,
        "改前",
        "Good day ... welcome to the CIFI Technologies Financial Results ...",
        accent=RED,
    )
    add_card(
        slide,
        0.75,
        3.55,
        5.9,
        1.28,
        "改后",
        "Good day ... welcome to the Sify Technologies Financial Results ...",
        accent=TEAL,
    )
    add_card(slide, 0.75, 5.05, 5.9, 0.9, "证据词", "Sify Technologies | non-GAAP | IFRS", accent=BLUE)
    add_bullets(
        slide,
        7.05,
        2.05,
        5.35,
        3.0,
        [
            "默认样本：Earnings-22 真实财报电话，多说话人。",
            "用户可以播放音频，也能看到 raw → corrected 对照。",
            "其他 AMI/AISHELL 样本保留为 speaker/overlap EvidenceMap 示例。",
        ],
        size=18,
    )
    add_textbox(slide, 7.05, 5.28, 5.35, 0.45, "https://fyangmie-talkweaver.hf.space", size=18, bold=True, color=TEAL)
    add_footer(slide, 5)

    # 6. Data and real evaluation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "真实公开数据验证", "主张来自仓库已有 CSV/文档；不把 controlled 结果当真实泛化")
    add_card(slide, 0.7, 2.05, 2.9, 1.45, "FLEURS", "30 clips\n中/英/法 read speech\n多语言 ASR sanity", accent=TEAL)
    add_card(slide, 3.82, 2.05, 2.9, 1.45, "AMI", "8 formal + 24 held-out\n英语会议\nASR + diarization + map", accent=BLUE)
    add_card(slide, 6.94, 2.05, 2.9, 1.45, "AISHELL-4", "60 × 20s\n普通话会议\nASR + DER/JER", accent=AMBER)
    add_card(slide, 10.06, 2.05, 2.55, 1.45, "Earnings-22", "财报电话\nRAG 专有词修正\nblind/smoke subsets", accent=RED)
    add_image(slide, "assets/result_charts/asr_error_by_dataset.png", 0.85, 4.0, 5.45, 2.25)
    add_image(slide, "assets/result_charts/asr_error_by_language.png", 6.95, 4.0, 5.45, 2.25)
    add_footer(slide, 6)

    # 7. Diarization and evidence maps
    slide = prs.slides.add_slide(blank)
    add_title(slide, "说话人/重叠证据：会议难点不只是文字", "pyannote 能提供有效 speaker-time evidence，但 overlap 仍然困难")
    add_metric(slide, 0.75, 2.05, 2.5, "AMI held-out", "DER 0.106", "JER 0.307", accent=TEAL)
    add_metric(slide, 3.55, 2.05, 2.5, "AMI overlap", "F1 0.490", "重叠仍难", accent=AMBER)
    add_metric(slide, 6.35, 2.05, 2.5, "AISHELL-4", "DER 0.327", "JER 0.713", accent=RED)
    add_metric(slide, 9.15, 2.05, 2.9, "Interruption", "10/10", "候选窗口人工确认", accent=BLUE)
    add_image(slide, "assets/result_charts/workflow_ablation_completeness.png", 0.9, 3.65, 5.4, 2.45)
    add_image(slide, "assets/result_charts/workflow_ablation_review_flags.png", 6.9, 3.65, 5.4, 2.45)
    add_footer(slide, 7)

    # 8. RAG term recovery
    slide = prs.slides.add_slide(blank)
    add_title(slide, "RAG 专有词恢复：有效，但必须保守", "贡献线不是泛化改写，而是 evidence-backed term rescue")
    add_card(slide, 0.75, 2.05, 3.65, 1.1, "典型修正", "CIFI Technologies → Sify Technologies\nnon-gap → non-GAAP\nIFR → IFRS", accent=TEAL)
    add_card(slide, 4.65, 2.05, 3.5, 1.1, "Earnings-22 v3", "base term recall\n0.833333 → 1.000000", accent=BLUE)
    add_card(slide, 8.4, 2.05, 3.8, 1.1, "边界", "WER 保持 0.212099 不变\n不是整体 ASR 提升", accent=AMBER)
    add_bullets(
        slide,
        0.9,
        3.7,
        5.65,
        2.15,
        [
            "v2 说明 RAG 可以帮助弱 ASR，但也会伤害强 ASR。",
            "v3 加安全门控：弱证据不改，常见词不乱映射为公司名。",
            "网页主 demo 用 Earnings-22 多说话人电话展示可播放术语修正。",
        ],
        size=17,
    )
    add_image(slide, "assets/result_charts/term_rescue_f1_by_variant.png", 7.0, 3.55, 5.2, 2.45)
    add_footer(slide, 8)

    # 9. Safety and limitations
    slide = prs.slides.add_slide(blank)
    add_title(slide, "安全审计与负结果", "论文价值来自诚实地区分正结果、controlled 结果和失败风险")
    add_image(slide, "assets/result_charts/evidence_gate_feature_set_macro_f1.png", 0.78, 2.05, 5.35, 2.35)
    add_image(slide, "assets/result_charts/evidence_gate_feature_set_unsafe_accept.png", 6.85, 2.05, 5.35, 2.35)
    add_bullets(
        slide,
        0.9,
        4.75,
        11.4,
        1.45,
        [
            "EvidenceGate 独立 held-out macro F1 仅约 0.325：不能主张已泛化。",
            "Controlled overlap safety 支持设计动机，但不能当真实会议性能。",
            "手机真机部署未完成；whisper.cpp 是 Level 1 local-machine benchmark。",
        ],
        size=17,
    )
    add_footer(slide, 9)

    # 10. Conclusion
    slide = prs.slides.add_slide(blank)
    add_title(slide, "结论与下一步", "TalkWeaver 的重点是可审计，不是宣称 SOTA")
    add_card(
        slide,
        0.8,
        2.05,
        3.7,
        2.0,
        "最强结论",
        "真实会议/电话转写需要 evidence map。TalkWeaver 能把 ASR、speaker、overlap、RAG、audit 组织成可检查结果。",
        accent=TEAL,
    )
    add_card(
        slide,
        4.85,
        2.05,
        3.7,
        2.0,
        "当前限制",
        "不是 SOTA；真实样本仍小；RAG 只稳定支持术语恢复；EvidenceGate 泛化仍弱。",
        accent=AMBER,
    )
    add_card(
        slide,
        8.9,
        2.05,
        3.45,
        2.0,
        "下一步",
        "扩大 held-out；增加真实人工标签；做真机移动 benchmark；完善论文和视频叙事。",
        accent=BLUE,
    )
    add_textbox(slide, 1.0, 4.85, 11.2, 0.5, "Demo: https://fyangmie-talkweaver.hf.space", size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.05, 5.62, 11.1, 0.45, "一句话：让混乱会议转写从“相信模型”变成“检查证据”。", size=21, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_footer(slide, 10)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    make_deck()
