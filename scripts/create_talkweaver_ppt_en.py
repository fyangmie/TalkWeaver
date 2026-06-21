#!/usr/bin/env python3
"""Generate an English TalkWeaver project overview deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from create_talkweaver_ppt import (
    AMBER,
    BLUE,
    INK,
    LINE,
    MUTED,
    PAPER,
    RED,
    ROOT,
    TEAL,
    WHITE,
    add_bullets,
    add_card,
    add_footer,
    add_image,
    add_metric,
    add_textbox,
)


OUT = ROOT / "docs" / "TalkWeaver_Project_Presentation_EN.pptx"


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, 0.62, 0.32, 9.9, 0.45, "TalkWeaver", size=13, bold=True, color=TEAL)
    add_textbox(slide, 0.62, 0.76, 10.8, 0.6, title, size=27, bold=True, color=INK)
    if subtitle:
        add_textbox(slide, 0.64, 1.35, 11.0, 0.42, subtitle, size=14, color=MUTED)
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.62),
        Inches(1.82),
        Inches(12.1),
        Inches(0.02),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def make_deck() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_textbox(slide, 0.72, 0.78, 11.4, 0.42, "AI Meeting Detective", size=16, bold=True, color=TEAL)
    add_textbox(slide, 0.72, 1.32, 11.5, 0.9, "TalkWeaver", size=52, bold=True, color=INK)
    add_textbox(
        slide,
        0.78,
        2.22,
        11.4,
        0.58,
        "Evidence-Grounded Conversation Maps for Chaotic Multi-Speaker Speech",
        size=23,
        bold=True,
        color=INK,
    )
    add_textbox(
        slide,
        0.8,
        3.12,
        10.7,
        0.8,
        "A research system for overlap, interruptions, misheard domain terms, speaker-time evidence, and correction audits.",
        size=18,
        color=MUTED,
    )
    add_card(slide, 0.82, 4.4, 3.45, 1.0, "Live Demo", "https://fyangmie-talkweaver.hf.space", accent=TEAL)
    add_card(slide, 4.55, 4.4, 3.45, 1.0, "Main Demo Clip", "Earnings-22 multi-speaker term correction", accent=BLUE)
    add_card(slide, 8.28, 4.4, 3.45, 1.0, "Core Framing", "Not a generic meeting summarizer; an evidence-auditing framework.", accent=AMBER)
    add_footer(slide, 1)

    # 2. Motivation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Motivation: Meeting Transcription Is Not Just WER", "Speakers, overlap, low-energy speech, domain terms, and LLM edits interact.")
    add_bullets(
        slide,
        0.75,
        2.18,
        5.7,
        3.45,
        [
            "A fluent transcript can still hide who spoke, when it was said, and whether a correction is justified.",
            "Chaotic meetings introduce speaker-boundary errors, overlap, interruptions, low-confidence speech, and domain-term misrecognition.",
            "LLM post-processing can make text look cleaner while silently adding unsupported edits.",
            "TalkWeaver treats meeting transcription as an evidence-grounded auditing problem.",
        ],
        size=17.5,
    )
    add_metric(slide, 7.0, 2.1, 2.7, "FLEURS zh-CN", "CER 0.113", "read speech, base", accent=TEAL)
    add_metric(slide, 10.0, 2.1, 2.7, "AISHELL-4", "CER 0.537", "meeting speech, base", accent=RED)
    add_metric(slide, 7.0, 3.55, 2.7, "AMI", "WER 0.398", "formal meeting, base", accent=AMBER)
    add_metric(slide, 10.0, 3.55, 2.7, "Takeaway", "Read speech != meetings", "evidence is needed", accent=BLUE)
    add_footer(slide, 2)

    # 3. RQs and contributions
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Research Questions and Contributions", "From automatic text rewriting to evidence-aware transcript auditing.")
    add_card(slide, 0.75, 2.1, 3.0, 1.35, "RQ1", "Can speaker-time structure make transcripts more auditable?", accent=TEAL)
    add_card(slide, 3.95, 2.1, 3.0, 1.35, "RQ2", "Can overlap-aware uncertainty reduce unsafe correction?", accent=AMBER)
    add_card(slide, 7.15, 2.1, 3.0, 1.35, "RQ3", "Can RAG glossary retrieval recover domain terms?", accent=BLUE)
    add_card(slide, 10.35, 2.1, 2.25, 1.35, "RQ4", "What is the local ASR speed/quality trade-off?", accent=RED)
    add_bullets(
        slide,
        0.92,
        4.1,
        11.4,
        1.8,
        [
            "Contribution 1: Temporal-anchor ConversationMap for ASR, speaker, time, overlap, retrieved terms, and audit trails.",
            "Contribution 2: Conservative RAG-based term recovery with evidence gates and review routing.",
            "Contribution 3: Safety-centered evaluation that separates real public data, controlled stress tests, demos, and negative results.",
        ],
        size=17.5,
    )
    add_footer(slide, 3)

    # 4. Architecture
    slide = prs.slides.add_slide(blank)
    add_title(slide, "System Architecture: From Audio to EvidenceMap", "Every stage preserves inputs, outputs, timestamps, and provenance.")
    add_image(slide, "assets/architecture.png", 0.8, 2.05, 5.45)
    add_bullets(
        slide,
        6.55,
        2.05,
        5.8,
        4.05,
        [
            "Audio preprocessing: standardize sample rate, channels, and volume.",
            "ASR: faster-whisper / whisper.cpp / deterministic mock fallback.",
            "Diarization: pyannote or reference speaker-time evidence.",
            "Overlap and event extraction: flag overlap, interruption candidates, and low-confidence regions.",
            "RAG term rescue: retrieve domain terms, company names, and financial terminology.",
            "Correction audit: preserve raw text, corrected text, evidence terms, and accept/reject/review decisions.",
        ],
        size=15.9,
    )
    add_footer(slide, 4)

    # 5. Web demo
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Web MVP: Show Exactly What Changed", "The deployed demo makes the correction evidence visible to non-expert users.")
    add_card(
        slide,
        0.75,
        2.05,
        5.9,
        1.28,
        "Before",
        "Good day ... welcome to the CIFI Technologies Financial Results ...",
        accent=RED,
    )
    add_card(
        slide,
        0.75,
        3.55,
        5.9,
        1.28,
        "After",
        "Good day ... welcome to the Sify Technologies Financial Results ...",
        accent=TEAL,
    )
    add_card(slide, 0.75, 5.05, 5.9, 0.9, "Evidence Terms", "Sify Technologies | non-GAAP | IFRS", accent=BLUE)
    add_bullets(
        slide,
        7.05,
        2.05,
        5.35,
        3.0,
        [
            "Default clip: real Earnings-22 financial-results call with multiple speakers.",
            "Users can play the audio and inspect raw-to-corrected changes.",
            "AMI/AISHELL examples remain available for speaker and overlap evidence.",
        ],
        size=18,
    )
    add_textbox(slide, 7.05, 5.28, 5.35, 0.45, "https://fyangmie-talkweaver.hf.space", size=18, bold=True, color=TEAL)
    add_footer(slide, 5)

    # 6. Datasets
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Real Public-Data Evaluation", "Reported claims come from existing CSVs and docs; controlled tests are not treated as real-world generalization.")
    add_card(slide, 0.7, 2.05, 2.9, 1.45, "FLEURS", "30 clips\nen/fr/zh-CN read speech\nmultilingual ASR sanity", accent=TEAL)
    add_card(slide, 3.82, 2.05, 2.9, 1.45, "AMI", "8 formal + 24 held-out\nEnglish meetings\nASR + diarization + maps", accent=BLUE)
    add_card(slide, 6.94, 2.05, 2.9, 1.45, "AISHELL-4", "60 × 20s\nMandarin meetings\nASR + DER/JER", accent=AMBER)
    add_card(slide, 10.06, 2.05, 2.55, 1.45, "Earnings-22", "financial calls\nRAG term correction\nblind/smoke subsets", accent=RED)
    add_image(slide, "assets/result_charts/asr_error_by_dataset.png", 0.85, 4.0, 5.45, 2.25)
    add_image(slide, "assets/result_charts/asr_error_by_language.png", 6.95, 4.0, 5.45, 2.25)
    add_footer(slide, 6)

    # 7. Diarization
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Speaker and Overlap Evidence", "Diarization is useful, but overlap remains a major source of uncertainty.")
    add_metric(slide, 0.75, 2.05, 2.5, "AMI held-out", "DER 0.106", "JER 0.307", accent=TEAL)
    add_metric(slide, 3.55, 2.05, 2.5, "AMI overlap", "F1 0.490", "still difficult", accent=AMBER)
    add_metric(slide, 6.35, 2.05, 2.5, "AISHELL-4", "DER 0.327", "JER 0.713", accent=RED)
    add_metric(slide, 9.15, 2.05, 2.9, "Interruption", "10/10", "candidate windows verified", accent=BLUE)
    add_image(slide, "assets/result_charts/workflow_ablation_completeness.png", 0.9, 3.65, 5.4, 2.45)
    add_image(slide, "assets/result_charts/workflow_ablation_review_flags.png", 6.9, 3.65, 5.4, 2.45)
    add_footer(slide, 7)

    # 8. RAG
    slide = prs.slides.add_slide(blank)
    add_title(slide, "RAG Term Recovery: Useful, but Must Be Conservative", "The contribution is evidence-backed term rescue, not unrestricted transcript rewriting.")
    add_card(slide, 0.75, 2.05, 3.65, 1.1, "Typical Fixes", "CIFI Technologies → Sify Technologies\nnon-gap → non-GAAP\nIFR → IFRS", accent=TEAL)
    add_card(slide, 4.65, 2.05, 3.5, 1.1, "Earnings-22 v3", "base term recall\n0.833333 → 1.000000", accent=BLUE)
    add_card(slide, 8.4, 2.05, 3.8, 1.1, "Boundary", "WER unchanged at 0.212099\nnot a general ASR improvement", accent=AMBER)
    add_bullets(
        slide,
        0.9,
        3.7,
        5.65,
        2.15,
        [
            "v2 showed that RAG can help weaker ASR but can harm stronger ASR.",
            "v3 adds safety gates: weak evidence is not applied; common words are not blindly mapped to company names.",
            "The web demo uses Earnings-22 to show playable, evidence-backed term correction.",
        ],
        size=16.8,
    )
    add_image(slide, "assets/result_charts/term_rescue_f1_by_variant.png", 7.0, 3.55, 5.2, 2.45)
    add_footer(slide, 8)

    # 9. Safety
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Safety Audits and Negative Results", "The paper should be honest about what works, what is controlled, and what fails.")
    add_image(slide, "assets/result_charts/evidence_gate_feature_set_macro_f1.png", 0.78, 2.05, 5.35, 2.35)
    add_image(slide, "assets/result_charts/evidence_gate_feature_set_unsafe_accept.png", 6.85, 2.05, 5.35, 2.35)
    add_bullets(
        slide,
        0.9,
        4.75,
        11.4,
        1.45,
        [
            "EvidenceGate independent held-out macro F1 is only about 0.325: no strong generalization claim.",
            "Controlled overlap safety supports the design, but it is not real meeting performance.",
            "Mobile deployment is not complete; whisper.cpp is a Level 1 local-machine benchmark.",
        ],
        size=17,
    )
    add_footer(slide, 9)

    # 10. Conclusion
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Conclusion and Next Steps", "TalkWeaver is about auditability, not state-of-the-art claims.")
    add_card(
        slide,
        0.8,
        2.05,
        3.7,
        2.0,
        "Strongest Claim",
        "Real meeting and call transcription benefits from an EvidenceMap that links ASR, speakers, overlap, RAG, and audits.",
        accent=TEAL,
    )
    add_card(
        slide,
        4.85,
        2.05,
        3.7,
        2.0,
        "Current Limits",
        "Not SOTA; real subsets are small; RAG mainly supports term recovery; EvidenceGate generalization remains weak.",
        accent=AMBER,
    )
    add_card(
        slide,
        8.9,
        2.05,
        3.45,
        2.0,
        "Next Steps",
        "Larger held-out evaluation, human labels, true mobile benchmark, and paper/video polishing.",
        accent=BLUE,
    )
    add_textbox(slide, 1.0, 4.85, 11.2, 0.5, "Demo: https://fyangmie-talkweaver.hf.space", size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.05, 5.62, 11.1, 0.45, "One-line takeaway: move from trusting fluent transcripts to inspecting evidence.", size=21, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_footer(slide, 10)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    make_deck()
